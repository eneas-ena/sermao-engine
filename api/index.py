from __future__ import annotations

import base64
import json
from io import BytesIO
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler
from pathlib import Path
from urllib.parse import parse_qs, urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from api.analyze import analyze_submission
from api.export import generate_export


STATIC_ROOT = Path(__file__).resolve().parents[1] / "static"
STATIC_FILES = {
    "": ("index.html", "text/html; charset=utf-8"),
    "/": ("index.html", "text/html; charset=utf-8"),
    "/static/index.html": ("index.html", "text/html; charset=utf-8"),
    "/static/styles.css": ("styles.css", "text/css; charset=utf-8"),
    "/static/app.js": ("app.js", "application/javascript; charset=utf-8"),
    "/static/app-icon.png": ("app-icon.png", "image/png"),
    "/static/manifest.webmanifest": ("manifest.webmanifest", "application/manifest+json; charset=utf-8"),
    "/favicon.ico": ("app-icon.png", "image/png"),
    "/apple-touch-icon.png": ("app-icon.png", "image/png"),
    "/apple-touch-icon-precomposed.png": ("app-icon.png", "image/png"),
}

NAVY = "041A2E"
NAVY_2 = "082B48"
GOLD = "DEA94A"
CREAM = "F7F2E8"
WHITE = "FFFFFF"
INK = "0F1A2B"
MUTED = "546273"


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _shape(slide, kind, x, y, width, height, fill, line=None):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line or fill)
    return shape


def _text(slide, value, x, y, width, height, size=22, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.05)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(value or "")
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_header(slide, data, heading, title_size):
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.08, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 1.05, 13.333, 0.03, GOLD)
    _text(slide, heading.upper(), 0.55, 0.12, 11.8, 0.58, title_size, GOLD, True)
    _text(slide, str(data.get("titulo", "")).upper(), 0.56, 0.7, 11.8, 0.25, 11, WHITE, True)


def _add_footer(slide):
    _text(slide, "Infográfico de sermão | conteúdo e elementos editáveis", 0.56, 7.0, 8.5, 0.22, 9, MUTED)


def _add_text_slide(prs, data, heading, content, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(CREAM)
    _add_header(slide, data, heading, style["titulo_slide"])
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 1.6, 11.85, 4.72, WHITE, GOLD)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.75, 1.6, 0.12, 4.72, GOLD)
    text_size = min(style["corpo"] + 5, 30)
    if len(str(content or "")) > 600:
        text_size = min(text_size, 20)
    _text(slide, content, 1.28, 1.95, 10.7, 3.95, text_size, INK)
    _add_footer(slide)


def build_pptx(data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    raw_style = data.get("estilo") or {}
    style = {
        "titulo": max(36, min(72, int(raw_style.get("titulo", 54)))),
        "titulo_slide": max(26, min(48, int(raw_style.get("titulo_slide", 34)))),
        "corpo": max(16, min(32, int(raw_style.get("corpo", 22)))),
    }

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, GOLD)
    image_data = data.get("imagem_capa", "")
    has_image = isinstance(image_data, str) and image_data.startswith("data:image/")
    text_width = 7.1 if has_image else 11.8
    _text(slide, str(data.get("titulo", "")).upper(), 0.7, 0.7, text_width, 2.4, style["titulo"], WHITE, True)
    _text(slide, str(data.get("subtitulo", "")).upper(), 0.72, 3.05, text_width, 1.05, max(22, round(style["titulo"] * 0.48)), GOLD, True)
    if has_image:
        try:
            image = BytesIO(base64.b64decode(image_data.split(",", 1)[1]))
            slide.shapes.add_picture(image, Inches(8.25), Inches(0.62), Inches(4.5), Inches(5.7))
        except (ValueError, IndexError):
            pass
    bible = data.get("texto_biblico") or {}
    bible_width = 7.1 if has_image else 11.25
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 5.12, bible_width, 1.45, NAVY_2, GOLD)
    _text(slide, f'TEXTO: {bible.get("referencia", "")}'.upper(), 1.0, 5.25, bible_width - 0.6, 0.4, 20, GOLD, True)
    _text(slide, bible.get("texto", ""), 1.0, 5.7, bible_width - 0.6, 0.65, 16, WHITE)

    _add_text_slide(prs, data, "Introdução", data.get("introducao", ""), style)
    for card in data.get("cards") or []:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(NAVY)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 0.72, 11.68, 5.92, CREAM, GOLD)
        _shape(slide, MSO_SHAPE.OVAL, 1.28, 1.2, 1.04, 1.04, GOLD)
        _text(slide, card.get("numero", ""), 1.28, 1.22, 1.04, 0.9, 36, NAVY, True, PP_ALIGN.CENTER)
        _text(slide, str(card.get("titulo", "")).upper(), 2.7, 1.05, 9.1, 1.2, style["titulo_slide"], NAVY, True)
        if card.get("referencia"):
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.7, 2.32, 3.0, 0.43, NAVY)
            _text(slide, card["referencia"], 2.9, 2.32, 2.6, 0.4, 14, WHITE, True)
        summary = str(card.get("resumo", ""))
        summary_size = 18 if len(summary) > 220 else 20 if len(summary) > 150 else min(style["corpo"] + 2, 24)
        _text(slide, summary, 1.42, 2.9, 10.5, 2.15, summary_size, INK)
        if card.get("frase"):
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.42, 5.28, 10.5, 0.9, NAVY)
            _text(slide, f'“{card["frase"]}”', 1.8, 5.35, 9.75, 0.65, min(style["corpo"] + 3, 25), GOLD, True, PP_ALIGN.CENTER)

    _add_text_slide(prs, data, "Aplicação prática", data.get("aplicacao", ""), style)
    _add_text_slide(prs, data, "Desafio", data.get("desafio", ""), style)
    _add_text_slide(prs, data, "Apelo", data.get("apelo", ""), style)
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


class handler(BaseHTTPRequestHandler):
    def do_GET(self) -> None:
        path = urlparse(self.path).path
        static_file = STATIC_FILES.get(path)
        if static_file is None:
            self.send_error(HTTPStatus.NOT_FOUND)
            return
        filename, content_type = static_file
        self._send((STATIC_ROOT / filename).read_bytes(), content_type)

    def do_POST(self) -> None:
        path = urlparse(self.path).path.rstrip("/")
        if path == "/api/analyze":
            self._analyze()
            return
        if path == "/api/export":
            self._export()
            return
        if path == "/api/pptx":
            self._pptx()
            return
        self.send_error(HTTPStatus.NOT_FOUND)

    def _analyze(self) -> None:
        import cgi

        form = cgi.FieldStorage(
            fp=self.rfile,
            headers=self.headers,
            environ={
                "REQUEST_METHOD": "POST",
                "CONTENT_TYPE": self.headers.get("Content-Type", ""),
            },
        )
        text = form.getfirst("text", "").strip()
        file_item = form["docx"] if "docx" in form else None
        docx_bytes = file_item.file.read() if file_item is not None and getattr(file_item, "filename", "") else None
        data = json.dumps(analyze_submission(text, docx_bytes), ensure_ascii=False).encode("utf-8")
        self._send(data, "application/json; charset=utf-8")

    def _export(self) -> None:
        length = int(self.headers.get("Content-Length", "0"))
        payload = json.loads(self.rfile.read(length).decode("utf-8"))
        kind = parse_qs(urlparse(self.path).query).get("kind", [""])[0]
        try:
            data, content_type, filename = generate_export(payload, kind)
        except ValueError:
            self.send_error(HTTPStatus.BAD_REQUEST, "Tipo de arquivo inválido")
            return
        self._send(data, content_type, filename)

    def _pptx(self) -> None:
        length = int(self.headers.get("Content-Length", "0"))
        payload = json.loads(self.rfile.read(length).decode("utf-8"))
        self._send(
            build_pptx(payload),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "slides_editaveis.pptx",
        )

    def _send(self, data: bytes, content_type: str, filename: str | None = None) -> None:
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-store")
        if filename:
            self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
        self.end_headers()
        self.wfile.write(data)
